import subprocess
import os
from pathlib import Path
import pandas as pd
import asyncio
from concurrent.futures import ThreadPoolExecutor, as_completed
from langchain_community.llms import Ollama
from langchain_community.document_loaders import PyPDFDirectoryLoader
from langchain_community.embeddings import OllamaEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.schema.document import Document
from langchain_community.vectorstores import FAISS
from langchain.chains import RetrievalQA
from typing import List, Optional
from pptx import Presentation
import logging

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def check_ollama_models():
    """Check if required Ollama models are available."""
    try:
        output = subprocess.run(["ollama", "list"], capture_output=True, text=True)
        installed_models = output.stdout.lower()
        return "llama3.2-vision" in installed_models and "nomic-embed-text" in installed_models
    except Exception as e:
        logging.error(f"Error checking Ollama models: {str(e)}")
        return False

def initialize_models():
    """Initialize the LLM and embedding models if they exist."""
    if not check_ollama_models():
        print("Required models are missing. Please install them manually using:")
        print("ollama pull llama3.2-vision\nollama pull nomic-embed-text")
        return None, None
    
    try:
        llm_model = Ollama(model="llama3.2-vision")
        embedding_model = OllamaEmbeddings(model="nomic-embed-text")
        return llm_model, embedding_model
    except Exception as e:
        logging.error(f"Error initializing models: {str(e)}")
        return None, None

def get_embedding_function(embedding_model):
    """Return the embedding function."""
    return embedding_model

def split_document(documents: List[Document]):
    """Split documents into smaller chunks for processing."""
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=80)
    return text_splitter.split_documents(documents)

def extract_text_from_shape(shape) -> str:
    """Extract text from a PowerPoint shape, including nested shapes."""
    text = []
    try:
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = " ".join(run.text.strip() for run in paragraph.runs)
                if full_text:
                    text.append(full_text)
        
        if hasattr(shape, "table"):
            text.append(extract_text_from_table(shape.table))
            
        if hasattr(shape, "shapes"):
            for sub_shape in shape.shapes:
                sub_text = extract_text_from_shape(sub_shape)
                if sub_text:
                    text.append(sub_text)
    except Exception as e:
        logging.warning(f"Error extracting text from shape: {str(e)}")
        
    return "\n".join(filter(None, text))

def extract_text_from_table(table) -> str:
    """Extract text from a PowerPoint table with improved formatting."""
    try:
        rows = []
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                rows.append(" | ".join(row_text))
        return "\n".join(rows)
    except Exception as e:
        logging.warning(f"Error extracting text from table: {str(e)}")
        return ""

def process_pptx_file(file_path: Path) -> List[Document]:
    """Process a PPTX file and extract all content with slide numbers, returning a list of Documents."""
    try:
        presentation = Presentation(file_path)
        documents = []
        
        for idx, slide in enumerate(presentation.slides, 1):
            slide_content = []
            
            # Extract content from all shapes in the slide
            for shape in slide.shapes:
                shape_text = extract_text_from_shape(shape)
                if shape_text:
                    slide_content.append(shape_text)
            
            if slide_content:  # Only create documents for slides with content
                full_content = f"=== Slide {idx} ===\n" + "\n".join(slide_content)
                documents.append(
                    Document(
                        page_content=full_content,
                        metadata={
                            "source": str(file_path),
                            "slide_number": idx
                        }
                    )
                )
        
        return documents
        
    except Exception as e:
        logging.error(f"Error processing PPTX file {file_path}: {str(e)}")
        return []

def convert_ppt_to_pptx(ppt_file: Path) -> Optional[Path]:
    """Convert PPT to PPTX using LibreOffice."""
    pptx_file = ppt_file.with_suffix(".pptx")
    try:
        subprocess.run([
            "soffice",
            "--headless",
            "--convert-to",
            "pptx",
            str(ppt_file),
            "--outdir",
            str(ppt_file.parent)
        ], check=True, capture_output=True)
        
        if pptx_file.exists():
            return pptx_file
            
    except subprocess.CalledProcessError as e:
        logging.error(f"Error converting PPT to PPTX: {str(e)}")
        return None
        
    return None

def load_text_files(directory_path):
    """Load text files from a directory."""
    documents = []
    for txt_file in Path(directory_path).glob("*.txt"):
        try:
            with open(txt_file, "r", encoding="utf-8") as file:
                content = file.read()
                if content.strip():
                    documents.append(Document(page_content=content, metadata={"source": str(txt_file)}))
        except Exception as e:
            logging.error(f"Error reading {txt_file}: {e}")
    return documents

def load_csv_files(directory_path):
    """Load CSV files from a directory."""
    documents = []
    for csv_file in Path(directory_path).glob("*.csv"):
        try:
            df = pd.read_csv(csv_file)
            for _, row in df.iterrows():
                content = row.to_string(index=False)
                if content.strip():
                    documents.append(Document(page_content=content, metadata={"source": str(csv_file)}))
        except Exception as e:
            logging.error(f"Error reading {csv_file}: {e}")
    return documents

def load_ppt_files(directory_path):
    """Load PowerPoint files from a directory."""
    documents = []
    
    # Handle PPTX files
    for pptx_file in Path(directory_path).glob("*.pptx"):
        try:
            slide_documents = process_pptx_file(pptx_file)
            documents.extend(slide_documents)
        except Exception as e:
            logging.error(f"Error processing PPTX file {pptx_file}: {e}")
    
    # Handle PPT files
    for ppt_file in Path(directory_path).glob("*.ppt"):
        try:
            converted_file = convert_ppt_to_pptx(ppt_file)
            if converted_file:
                slide_documents = process_pptx_file(converted_file)
                documents.extend(slide_documents)
                converted_file.unlink()  # Clean up temporary file
        except Exception as e:
            logging.error(f"Error processing PPT file {ppt_file}: {e}")
    
    return documents

async def load_documents_async(directory_path: str):
    """Asynchronously load documents from a directory."""
    if not os.path.isdir(directory_path):
        logging.error(f"Directory not found: {directory_path}")
        return []
    
    try:
        documents = []
        
        # Load PDF documents
        pdf_loader = PyPDFDirectoryLoader(directory_path)
        documents.extend(pdf_loader.load())
        
        # Load text files
        documents.extend(load_text_files(directory_path))
        
        # Load CSV files
        documents.extend(load_csv_files(directory_path))
        
        # Load PowerPoint files - now each slide becomes a separate document
        documents.extend(load_ppt_files(directory_path))
        
        # Split all documents into chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=80)
        chunks = text_splitter.split_documents(documents)
        
        logging.info(f"Loaded {len(documents)} documents and created {len(chunks)} chunks")
        return chunks
    except Exception as e:
        logging.error(f"Error loading documents: {str(e)}")
        return []
    
    
def add_to_faiss(chunks: List[Document], save_path: str, embedding_model):
    """Add document chunks to FAISS index."""
    try:
        db = FAISS.from_documents(chunks, get_embedding_function(embedding_model))
        os.makedirs(os.path.dirname(save_path), exist_ok=True)
        db.save_local(save_path)
        return db
    except Exception as e:
        logging.error(f"Error adding to FAISS: {str(e)}")
        return None

def load_or_create_faiss(directory_path: str, faiss_index_path: str, embedding_model):
    """Load or create a FAISS index."""
    index_path = Path(faiss_index_path)
    try:
        if (index_path / "index.faiss").exists() and (index_path / "index.pkl").exists():
            return FAISS.load_local(str(index_path), get_embedding_function(embedding_model), allow_dangerous_deserialization=True)
        
        documents = asyncio.run(load_documents_async(directory_path))
        if documents:
            chunks = split_document(documents)
            return add_to_faiss(chunks, str(index_path), embedding_model) if chunks else None
    except Exception as e:
        logging.error(f"Error loading/creating FAISS index: {str(e)}")
    return None

def create_qa_chain(db, llm_model):
    """Create a Retrieval QA chain that includes source documents."""
    try:
        retriever = db.as_retriever(search_type="similarity", search_kwargs={"k": 5})
        return RetrievalQA.from_chain_type(
            llm=llm_model, 
            chain_type="stuff", 
            retriever=retriever,
            return_source_documents=True  # Enable source document return
        )
    except Exception as e:
        logging.error(f"Error creating QA chain: {str(e)}")
        return None

def print_chunks(db):
    """Print all chunks stored in the FAISS database."""
    try:
        # Get all document chunks from the FAISS index
        all_docs = db.docstore._dict
        
        print("\n=== ALL CHUNKS ===\n")
        for i, (doc_id, doc) in enumerate(all_docs.items(), 1):
            print(f"\nCHUNK {i}:")
            print(f"Document ID: {doc_id}")
            print(f"Source: {doc.metadata.get('source', 'Unknown')}")
            if 'slide_number' in doc.metadata:
                print(f"Slide Number: {doc.metadata['slide_number']}")
            print("\nContent:")
            print("-" * 50)
            print(doc.page_content)
            print("-" * 50)
            print()
        
        print(f"\nTotal number of chunks: {len(all_docs)}")
        
    except Exception as e:
        logging.error(f"Error printing chunks: {str(e)}")
        print("Failed to print chunks.")

def get_unique_sources(source_documents):
    """Extract unique sources from the source documents."""
    unique_sources = set()
    for doc in source_documents:
        source = doc.metadata.get('source', 'Unknown')
        slide_number = doc.metadata.get('slide_number')
        if slide_number is not None:
            source = f"{source} (Slide {slide_number})"
        unique_sources.add(source)
    return sorted(list(unique_sources))

def get_response(qa_chain, query: str, db):
    """Get response from the QA chain and return the result and sources."""
    try:
        if query.lower() == "chunk":
            print_chunks(db)
            return "Displayed all chunks in the console.", []
        else:
            result = qa_chain.invoke(query)
            response = result['result']
            
            # Get unique sources
            sources = []
            if 'source_documents' in result:
                sources = get_unique_sources(result['source_documents'])
            
            return response, sources
            
    except Exception as e:
        logging.error(f"Error getting response: {str(e)}")
        return "Failed to get response. Please try another question.", []
        
def main():
    """Main entry point for the program."""
    llm_model, embedding_model = initialize_models()
    if not llm_model or not embedding_model:
        print("Failed to initialize models. Exiting...")
        return

    directory_path = input("Enter the path to your document directory: ").strip()
    if not os.path.isdir(directory_path):
        print("Invalid directory path. Exiting...")
        return

    print("\nLoading and processing documents from the directory...")
    faiss_index_path = os.path.join(directory_path, "faiss_index")
    db = load_or_create_faiss(directory_path, faiss_index_path, embedding_model)
    if not db:
        print("Failed to initialize the database.")
        return

    qa_chain = create_qa_chain(db, llm_model)
    if not qa_chain:
        print("Failed to create QA chain.")
        return

    print("\nEnter your questions (type 'exit' to quit, 'chunk' to see all chunks):")
    while True:
        query = input("\nYour question: ").strip()
        if query.lower() == 'exit':
            break
        if query:
            get_response(qa_chain, query, db)
        else:
            print("Please enter a valid question.")

if __name__ == "__main__":
    main()